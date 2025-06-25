import os
import sys
import zipfile
from fpdf import FPDF
from PIL import Image
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader
import fitz  # PyMuPDF


# -------- PDF Conversion -------- #
def to_pdf(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    output_path = os.path.splitext(file_path)[0] + ".pdf"

    try:
        if ext == ".txt":
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    pdf.cell(0, 10, txt=line.strip(), ln=True)
            pdf.output(output_path)

        elif ext in [".jpg", ".jpeg", ".png"]:
            image = Image.open(file_path)
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(output_path, "PDF")

        elif ext == ".docx":
            from docx2pdf import convert
            convert(file_path)
            output_path = os.path.splitext(file_path)[0] + ".pdf"

        elif ext == ".pptx":
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
            output_path = os.path.abspath(output_path)
            presentation.SaveAs(output_path, 32)  # 32 is PDF format
            presentation.Close()
            powerpoint.Quit()

        else:
            print("❌ Unsupported format for PDF conversion.")
            return

        print(f"✅ Converted to PDF: {output_path}")
    except Exception as e:
        print(f"❌ Error in to_pdf: {e}")


# -------- DOCX Conversion -------- #
def to_docx(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    output_path = os.path.splitext(file_path)[0] + ".docx"

    try:
        doc = Document()

        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    doc.add_paragraph(line.strip())

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        image = shape.image
                        img_data = image.blob
                        img_ext = image.ext
                        img_name = f"slide_img_{i}.{img_ext}"
                        with open(img_name, "wb") as f:
                            f.write(img_data)
                        doc.add_paragraph().add_run().add_picture(img_name, width=Inches(5))
                        os.remove(img_name)
                    elif hasattr(shape, "text") and shape.text.strip():
                        doc.add_paragraph(shape.text.strip())

        elif ext == ".pdf":
            pdf = fitz.open(file_path)
            for i, page in enumerate(pdf):
                pix = page.get_pixmap(dpi=150)
                img_path = f"pdf_page_{i}.png"
                pix.save(img_path)
                doc.add_paragraph().add_run().add_picture(img_path, width=Inches(6))
                os.remove(img_path)

        elif ext in [".jpg", ".jpeg", ".png"]:
            image = Image.open(file_path)
            width, height = image.size
            dpi = 96
            max_width = 6
            image_width_inches = width / dpi

            if image_width_inches > max_width:
                image_width_inches = max_width

            paragraph = doc.add_paragraph()
            run = paragraph.add_run()
            run.add_picture(file_path, width=Inches(image_width_inches))
            paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        elif ext == ".docx":
            print("✅ Already a DOCX file.")
            return

        else:
            print("❌ Unsupported format for DOCX conversion.")
            return

        doc.save(output_path)
        print(f"✅ Converted to DOCX: {output_path}")

    except Exception as e:
        print(f"❌ Error in to_docx: {e}")


# -------- PPTX Conversion -------- #
def to_pptx(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    output_path = os.path.splitext(file_path)[0] + ".pptx"

    try:
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        prs_width = prs.slide_width
        prs_height = prs.slide_height

        if ext == ".pdf":
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)
                image_path = f"temp_page_{page_num}.png"
                pix.save(image_path)

                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(image_path, 0, 0, width=prs_width, height=prs_height)
                os.remove(image_path)

        elif ext in [".jpg", ".jpeg", ".png"]:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(file_path, 0, 0, width=prs_width, height=prs_height)

        elif ext == ".docx":
            doc = Document(file_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            slides = [paragraphs[i:i + 7] for i in range(0, len(paragraphs), 7)]

            for slide_text in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                body = slide.shapes.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for line in slide_text:
                    tf.add_paragraph().text = line

        else:
            print("❌ Only PDF, DOCX, and image formats supported.")
            return

        prs.save(output_path)
        print(f"✅ Converted to PPTX: {output_path}")

    except Exception as e:
        print(f"❌ Error in to_pptx: {e}")


# -------- JPG Conversion -------- #
def to_jpg(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    output_path = os.path.splitext(file_path)[0] + ".jpg"

    try:
        if ext in [".png", ".jpeg", ".jpg"]:
            img = Image.open(file_path)
            rgb_im = img.convert("RGB")
            rgb_im.save(output_path, quality=95)

        elif ext == ".pdf":
            from pdf2image import convert_from_path
            images = convert_from_path(file_path)
            images[0].save(output_path, "JPEG")

        else:
            print("❌ JPG conversion only supports image or PDF files.")
            return

        print(f"✅ Converted to JPG: {output_path}")
    except Exception as e:
        print(f"❌ Error in to_jpg: {e}")


# -------- ZIP Compression -------- #
def compress_file(file_path):
    if not os.path.isfile(file_path):
        print("❌ File does not exist:", file_path)
        return
    try:
        zip_path = os.path.splitext(file_path)[0] + ".zip"
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            zipf.write(file_path, os.path.basename(file_path))
        print(f"✅ Compressed to ZIP: {zip_path}")
    except Exception as e:
        print(f"❌ Error in compress_file: {e}")
